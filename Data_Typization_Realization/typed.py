import zipfile
import io
import json
import os
import tempfile
from pathlib import Path
from datetime import datetime
from typing import Dict, List

# Библиотеки для документов
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
import rarfile


class DocumentTypizer:
    """
    Универсальный типизатор документов.
    Поддерживает: PDF, DOCX, DOC, PPTX, XLSX, XLS, ZIP, RAR, TXT, CSV, JSON, XML
    
    Пример использования:
        typizer = DocumentTypizer(output_dir=r"C:\Results")
        typizer.process_zip_file(r"C:\archive.zip")
    """
    
    def __init__(self, output_dir: str = None, max_text_size: int = None):
        """
        Args:
            output_dir: Папка для сохранения результатов (по умолчанию - папка 'results' рядом со скриптом)
            max_text_size: Максимальный размер текстовых файлов в символах (None = без ограничений)
        """

        if output_dir is None:
            self.output_dir = os.path.join(os.getcwd(), "results")
        else:
            self.output_dir = output_dir
        

        os.makedirs(self.output_dir, exist_ok=True)
        

        self.max_text_size = max_text_size
        

        self.processed_count = 0
        self.error_count = 0
        self.skipped_count = 0
        

        if os.name == 'nt':  
            rarfile.UNRAR_TOOL = r"C:\Program Files\WinRAR\UnRAR.exe"
        else:  
            rarfile.UNRAR_TOOL = "unrar"
        
        print(f"Папка сохранения: {self.output_dir}")
    

    def process_zip_file(self, zip_path: str, output_name: str = None) -> List[Dict]:
        """
        Главный метод - обработка ZIP файла
        
        Args:
            zip_path: Путь к ZIP файлу
            output_name: Имя выходного JSON файла (если None - генерируется автоматически)
        
        Returns:
            Список обработанных документов
        """
        print("=" * 70)
        print("ТИПИЗАТОР ДОКУМЕНТОВ v2.0")
        print("=" * 70)
        print(f"Архив: {zip_path}")
        

        if output_name is None:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            zip_name = Path(zip_path).stem
            output_name = f"{zip_name}_{timestamp}.json"
        
        output_path = os.path.join(self.output_dir, output_name)
        print(f"Результат: {output_path}")
        print("=" * 70)
        

        if not os.path.exists(zip_path):
            print(f"ОШИБКА: Файл не найден: {zip_path}")
            return []
        
        with open(zip_path, 'rb') as f:
            archive_data = f.read()
        
        print(f"Размер архива: {len(archive_data) / (1024*1024):.1f} МБ")
        
        documents = self._process_archive(archive_data, Path(zip_path).stem, "")
        

        self._save_to_json(documents, output_path)
        
        self._print_stats(documents)
        
        return documents
    
    
    def _process_archive(self, archive_data: bytes, archive_name: str, parent_path: str) -> List[Dict]:
        """Рекурсивная обработка ZIP архива"""
        documents = []
        
        try:
            with zipfile.ZipFile(io.BytesIO(archive_data)) as zip_ref:
                items = [f for f in zip_ref.infolist() if not f.is_dir()]
                
                print(f"\nАрхив: {archive_name} ({len(items)} файлов)")
                
                for i, file_info in enumerate(items, 1):
                    file_path = self._build_path(parent_path, archive_name, file_info.filename)
                    ext = Path(file_info.filename).suffix.lower()
                    
                    print(f"  [{i}/{len(items)}] {file_path}")
                    
                    try:
                        file_data = zip_ref.read(file_info)
                        doc = self._process_file(file_data, file_path, ext)
                        
                        if doc:
                            documents.append(doc)
                            self._update_counters(doc)
                    
                    except Exception as e:
                        print(f"    ОШИБКА: {e}")
                        documents.append(self._error_doc(file_path, ext, str(e)))
                        self.error_count += 1
        
        except zipfile.BadZipFile:
            print(f"  ОШИБКА: Не является ZIP архивом")
            documents.append(self._error_doc(archive_name, '.zip', 'Невалидный ZIP архив'))
            self.error_count += 1
        
        return documents
    
    def _process_rar(self, file_data: bytes, file_path: str) -> Dict:
        """Обработка RAR архива"""
        tmp_path = None
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.rar', delete=False) as tmp:
                tmp.write(file_data)
                tmp_path = tmp.name
            
            with rarfile.RarFile(tmp_path) as rf:
                files_list = []
                dirs_list = []
                total_size = 0
                
                for info in rf.infolist():
                    if info.is_dir():
                        dirs_list.append(info.filename)
                    else:
                        files_list.append({
                            "имя": info.filename,
                            "размер_КБ": round(info.file_size / 1024, 1)
                        })
                        total_size += info.file_size
                
                content = []
                content.append(f"RAR архив: {len(files_list)} файлов, {total_size/1024:.1f} КБ")
                
                if dirs_list:
                    content.append(f"\nПапки ({len(dirs_list)}):")
                    for d in sorted(dirs_list)[:20]:
                        content.append(f"  {d}")
                
                content.append(f"\nФайлы:")
                for f in files_list[:50]:
                    content.append(f"  {f['имя']} ({f['размер_КБ']} КБ)")
                
                if len(files_list) > 50:
                    content.append(f"  ... и еще {len(files_list) - 50} файлов")
                
                print(f"    Успешно")
                
                return {
                    "название": file_path,
                    "содержание": "\n".join(content),
                    "тип_файла": ".rar",
                    "количество_файлов": len(files_list),
                    "количество_папок": len(dirs_list),
                    "размер_КБ": round(total_size / 1024, 1),
                    "список_файлов": [f["имя"] for f in files_list[:100]],
                    "статус": "success"
                }
        
        except Exception as e:
            print(f"    ОШИБКА: {e}")
            return self._error_doc(file_path, '.rar', str(e))
        
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except:
                    pass
    
    
    def _process_pdf(self, file_data: bytes, file_path: str) -> Dict:
        """Извлечение текста из PDF"""
        try:
            pdf_file = io.BytesIO(file_data)
            reader = PdfReader(pdf_file)
            
            pages_text = []
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages_text.append(f"[Страница {page_num}]\n{text}")
            
            content = "\n\n".join(pages_text) if pages_text else "PDF без текстового слоя"
            
            print(f"    {len(reader.pages)} стр.")
            
            return {
                "название": file_path,
                "содержание": content,
                "тип_файла": ".pdf",
                "количество_страниц": len(reader.pages),
                "статус": "success"
            }
        
        except Exception as e:
            print(f"    ОШИБКА: {e}")
            return self._error_doc(file_path, '.pdf', str(e))
    

    def _process_docx(self, file_data: bytes, file_path: str) -> Dict:
        """Извлечение текста из Word"""
        try:
            doc_file = io.BytesIO(file_data)
            doc = Document(doc_file)
            
            parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            
            tables_count = len(doc.tables)
            for t_num, table in enumerate(doc.tables, 1):
                parts.append(f"\n[Таблица {t_num}]")
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    parts.append(" | ".join(row_data))
            
            content = "\n".join(parts) if parts else "Документ пуст"
            
            print(f"    {len(doc.paragraphs)} параграфов, {tables_count} таблиц")
            
            return {
                "название": file_path,
                "содержание": content,
                "тип_файла": Path(file_path).suffix.lower(),
                "количество_параграфов": len(doc.paragraphs),
                "количество_таблиц": tables_count,
                "статус": "success"
            }
        
        except Exception as e:
            print(f"    ОШИБКА: {e}")
            return self._error_doc(file_path, Path(file_path).suffix.lower(), str(e))
    

    
    def _process_pptx(self, file_data: bytes, file_path: str) -> Dict:
        """Извлечение текста из PowerPoint"""
        try:
            pptx_file = io.BytesIO(file_data)
            prs = Presentation(pptx_file)
            
            slides_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"[Слайд {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                    
                    if shape.has_table:
                        table_parts = ["\nТаблица:"]
                        for row in shape.table.rows:
                            cells = [cell.text for cell in row.cells]
                            table_parts.append(" | ".join(cells))
                        slide_content.extend(table_parts)
                
                slides_text.append("\n".join(slide_content))
            
            content = "\n\n".join(slides_text) if slides_text else "Презентация пуста"
            
            print(f"    {len(prs.slides)} слайдов")
            
            return {
                "название": file_path,
                "содержание": content,
                "тип_файла": ".pptx",
                "количество_слайдов": len(prs.slides),
                "статус": "success"
            }
        
        except Exception as e:
            print(f"    ОШИБКА: {e}")
            return self._error_doc(file_path, '.pptx', str(e))
    
    
    def _process_xlsx(self, file_data: bytes, file_path: str) -> Dict:
        """Извлечение данных из Excel"""
        try:
            excel_file = io.BytesIO(file_data)
            workbook = openpyxl.load_workbook(excel_file, data_only=True)
            
            all_sheets = []
            total_rows = 0
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                max_row = sheet.max_row
                max_col = sheet.max_column
                
                sheet_content = [f"[Лист: {sheet_name}] ({max_row} строк, {max_col} столбцов)", ""]
                
                row_count = 0
                for row in sheet.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):
                        sheet_content.append(" | ".join(row_data))
                        row_count += 1
                
                all_sheets.append("\n".join(sheet_content))
                total_rows += row_count
            
            separator = "\n\n" + "=" * 50 + "\n\n"
            content = separator.join(all_sheets)
            
            print(f"    {len(workbook.sheetnames)} листов, {total_rows} строк")
            
            return {
                "название": file_path,
                "содержание": content,
                "тип_файла": Path(file_path).suffix.lower(),
                "количество_листов": len(workbook.sheetnames),
                "всего_строк": total_rows,
                "листы": workbook.sheetnames,
                "статус": "success"
            }
        
        except Exception as e:
            print(f"    ОШИБКА: {e}")
            return self._error_doc(file_path, Path(file_path).suffix.lower(), str(e))
    
    
    def _process_text(self, file_data: bytes, file_path: str) -> Dict:
        """Обработка текстовых файлов"""
        ext = Path(file_path).suffix.lower()
        

        for encoding in ['utf-8', 'cp1251', 'latin-1']:
            try:
                text = file_data.decode(encoding)
                break
            except:
                continue
        else:
            text = file_data.decode('utf-8', errors='ignore')
        

        original_size = len(text)
        if self.max_text_size and len(text) > self.max_text_size:
            text = text[:self.max_text_size] + f"\n\n... [текст обрезан до {self.max_text_size} символов, всего было {original_size}]"
            print(f"    Обрезан до {self.max_text_size} символов (было {original_size})")
        else:
            print(f"    {original_size} символов")
        
        return {
            "название": file_path,
            "содержание": text,
            "тип_файла": ext,
            "размер_символов": original_size,
            "статус": "success"
        }
    
    
    def _process_file(self, file_data: bytes, file_path: str, ext: str) -> Dict:
        """Маршрутизатор по типам файлов"""
        
        if ext == '.zip':
            print(f"    Вложенный архив")

            self._process_archive(file_data, Path(file_path).stem, str(Path(file_path).parent))

            return self._get_zip_info(file_data, file_path)
        
        elif ext == '.rar':
            return self._process_rar(file_data, file_path)
        

        elif ext == '.pdf':
            return self._process_pdf(file_data, file_path)
        
        elif ext in ['.docx', '.doc']:
            return self._process_docx(file_data, file_path)
        
        elif ext == '.pptx':
            return self._process_pptx(file_data, file_path)
        
        elif ext in ['.xlsx', '.xls']:
            return self._process_xlsx(file_data, file_path)
        

        elif ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.py', '.log', '.ini', '.cfg']:
            return self._process_text(file_data, file_path)

        else:
            print(f"    Неподдерживаемый формат")
            self.skipped_count += 1
            return {
                "название": file_path,
                "содержание": f"Неподдерживаемый формат: {ext}",
                "тип_файла": ext,
                "статус": "unsupported"
            }
    

    
    def _get_zip_info(self, zip_data: bytes, file_path: str) -> Dict:
        """Информация о ZIP архиве"""
        try:
            with zipfile.ZipFile(io.BytesIO(zip_data)) as zf:
                files = [f.filename for f in zf.infolist() if not f.is_dir()]
                total_size = sum(f.file_size for f in zf.infolist())
                
                content = [f"ZIP архив: {len(files)} файлов, {total_size/1024:.1f} КБ"]
                content.append("\nСодержимое:")
                for f in files[:50]:
                    content.append(f"  {f}")
                
                if len(files) > 50:
                    content.append(f"  ... и еще {len(files) - 50} файлов")
                
                return {
                    "название": file_path,
                    "содержание": "\n".join(content),
                    "тип_файла": ".zip",
                    "количество_файлов": len(files),
                    "размер_КБ": round(total_size / 1024, 1),
                    "список_файлов": files[:100],
                    "статус": "success"
                }
        except:
            return self._error_doc(file_path, '.zip', 'Ошибка чтения ZIP')
    
    def _build_path(self, parent: str, archive: str, filename: str) -> str:
        """Построение полного пути с сохранением структуры папок"""
        if parent:
            path = f"{parent}/{archive}/{filename}"
        else:
            path = f"{archive}/{filename}"
        return path.replace('//', '/')
    
    def _error_doc(self, path: str, ext: str, error: str) -> Dict:
        """Создание документа с ошибкой"""
        return {
            "название": path,
            "содержание": f"ОШИБКА: {error}",
            "тип_файла": ext,
            "статус": "error"
        }
    
    def _update_counters(self, doc: Dict):
        """Обновление счетчиков"""
        status = doc.get("статус")
        if status == "success":
            self.processed_count += 1
        elif status == "error":
            self.error_count += 1
        else:
            self.skipped_count += 1
    

    def _save_to_json(self, documents: List[Dict], output_path: str):
        output = {
            "метаданные": {
                "дата_обработки": datetime.now().isoformat(),
                "файл_результата": os.path.basename(output_path),
                "статистика": {
                    "всего_документов": len(documents),
                    "успешно": self.processed_count,
                    "с_ошибками": self.error_count,
                    "пропущено": self.skipped_count
                },
                "типы_файлов": self._get_types_stats(documents)
            },
            "документы": documents
        }
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(output, f, ensure_ascii=False, indent=2)
        
        file_size = os.path.getsize(output_path)
        print(f"\nРезультат сохранен: {output_path}")
        print(f"Размер файла: {file_size / (1024*1024):.2f} МБ")
    
    def _get_types_stats(self, documents: List[Dict]) -> Dict:
        stats = {}
        for doc in documents:
            ext = doc.get("тип_файла", "неизвестный")
            status = doc.get("статус", "unknown")
            
            if ext not in stats:
                stats[ext] = {"всего": 0, "успешно": 0, "ошибок": 0, "пропущено": 0}
            
            stats[ext]["всего"] += 1
            if status == "success":
                stats[ext]["успешно"] += 1
            elif status == "error":
                stats[ext]["ошибок"] += 1
            else:
                stats[ext]["пропущено"] += 1
        
        return stats
    
    def _print_stats(self, documents: List[Dict]):
        print("\n" + "=" * 70)
        print("СТАТИСТИКА ОБРАБОТКИ")
        print("=" * 70)
        
        stats = self._get_types_stats(documents)
        for ext, stat in sorted(stats.items()):
            print(f"  {ext}: {stat['всего']} всего", end="")
            details = []
            if stat['успешно']:
                details.append(f"успешно: {stat['успешно']}")
            if stat['ошибок']:
                details.append(f"ошибок: {stat['ошибок']}")
            if stat['пропущено']:
                details.append(f"пропущено: {stat['пропущено']}")
            if details:
                print(f" ({', '.join(details)})")
            else:
                print()
        
        print(f"\n  Всего документов: {len(documents)}")
        print(f"  Успешно: {self.processed_count}")
        print(f"  С ошибками: {self.error_count}")
        print(f"  Пропущено: {self.skipped_count}")
        print("=" * 70)



# ============================================

if __name__ == "__main__":
    
    ZIP_FILE = r"D:\Хрень\Новая папка.zip"
    
    SAVE_DIR = r"C:\Users\bublick\Documents"
    
    typizer = DocumentTypizer(
        output_dir=SAVE_DIR,     
        max_text_size=None        
    )
    
    try:
        documents = typizer.process_zip_file(ZIP_FILE)
        
        if documents:
            print("\n" + "=" * 70)
            print("ПРИМЕРЫ ОБРАБОТАННЫХ ДОКУМЕНТОВ")
            print("=" * 70)
            
            for i, doc in enumerate(documents[:5], 1):
                print(f"\n{i}. {doc['название']}")
                print(f"   Тип: {doc['тип_файла']} | Статус: {doc['статус']}")
                
                if doc['статус'] == 'success':
                    preview = doc['содержание'][:200].replace('\n', ' ')
                    print(f"   Превью: {preview}...")
                else:
                    print(f"   {doc['содержание']}")
    
    except FileNotFoundError:
        print(f"ОШИБКА: Файл не найден: {ZIP_FILE}")
        print("Проверьте путь к ZIP архиву")
    
    except Exception as e:
        print(f"КРИТИЧЕСКАЯ ОШИБКА: {e}")
        import traceback
        traceback.print_exc()
    
    finally:
        print("\n" + "=" * 70)
        print("Программа завершена")
        print("=" * 70)